"""
file_generator.py
Pure functions to generate Office files from structured data.
No FastAPI or DB dependencies — importable standalone for testing.
"""

import re
import uuid
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────
# DOCX
# ──────────────────────────────────────────────

def generate_docx(title: str, content: str, save_path: Path) -> Path:
    """
    Create a Word document from a markdown-like string.
    Supports: ## heading2, ### heading3, - bullet, ** bold **, plain paragraphs.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title_para = doc.add_heading(title, level=1)
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

    for line in content.splitlines():
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            _add_inline_formatting(p, stripped[2:])
        elif re.match(r"^\d+\.\s", stripped):
            p = doc.add_paragraph(style="List Number")
            _add_inline_formatting(p, re.sub(r"^\d+\.\s", "", stripped))
        else:
            p = doc.add_paragraph()
            _add_inline_formatting(p, stripped)

    save_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(save_path))
    logger.info(f"[file_generator] Saved DOCX: {save_path}")
    return save_path


def _add_inline_formatting(paragraph, text: str):
    """Handles **bold** and *italic* within a paragraph run."""
    from docx.shared import Pt

    # Split on bold/italic markers
    parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith("*") and part.endswith("*"):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        else:
            paragraph.add_run(part)


# ──────────────────────────────────────────────
# XLSX
# ──────────────────────────────────────────────

def generate_xlsx(title: str, rows: list, save_path: Path) -> Path:
    """
    Create an Excel spreadsheet from a list-of-lists.
    First row is treated as the header (bold + blue fill).
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel sheet name max 31 chars

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center")

    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    col_widths: dict[int, int] = {}

    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_value in enumerate(row_data, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=str(cell_value) if cell_value is not None else "")
            cell.border = thin_border

            if row_idx == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
            else:
                cell.alignment = Alignment(vertical="center")

            # Track max col width
            cell_len = len(str(cell_value)) if cell_value is not None else 0
            col_widths[col_idx] = max(col_widths.get(col_idx, 0), cell_len)

    # Auto-size columns (capped at 50)
    for col_idx, width in col_widths.items():
        ws.column_dimensions[get_column_letter(col_idx)].width = min(width + 4, 50)

    # Freeze header row
    if rows:
        ws.freeze_panes = "A2"

    save_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(save_path))
    logger.info(f"[file_generator] Saved XLSX: {save_path}")
    return save_path


# ──────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────

def generate_pptx(title: str, slides: list, save_path: Path) -> Path:
    """
    Create a PowerPoint presentation.
    `slides` is a list of {"title": str, "bullets": list[str]}.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Title Slide ──────────────────────────────
    title_layout = prs.slide_layouts[0]  # "Title Slide"
    slide = prs.slides.add_slide(title_layout)

    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2F, 0x54, 0x96)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    # ── Content Slides ───────────────────────────
    bullet_layout = prs.slide_layouts[1]  # "Title and Content"

    for slide_data in slides:
        slide_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])

        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = slide_title
        s.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2F, 0x54, 0x96)
        s.shapes.title.text_frame.paragraphs[0].font.bold = True

        tf = s.placeholders[1].text_frame
        tf.clear()

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(bullet)
            p.level = 0
            p.font.size = Pt(20)

    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))
    logger.info(f"[file_generator] Saved PPTX: {save_path}")
    return save_path


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def make_safe_filename(text: str) -> str:
    """Turn an arbitrary title string into a filesystem-safe filename stem."""
    safe = re.sub(r"[^\w\s-]", "", text).strip()
    safe = re.sub(r"[\s_]+", "_", safe)
    return safe[:60] or "file"


def build_save_path(upload_dir: str, user_id: str, title: str, ext: str) -> Path:
    """Returns a unique save path inside uploads/generated/{user_id}/."""
    stem = make_safe_filename(title)
    filename = f"{uuid.uuid4().hex}_{stem}{ext}"
    return Path(upload_dir) / "generated" / str(user_id) / filename
